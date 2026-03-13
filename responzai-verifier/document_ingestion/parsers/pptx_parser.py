# document_ingestion/parsers/pptx_parser.py

from pptx import Presentation

async def parse_pptx(file_path: str) -> dict:
    """
    Extrahiert Text aus PowerPoint-Präsentationen.

    PowerPoint hat zwei Textquellen:
    1. Der sichtbare Folieninhalt (Textboxen, Titel)
    2. Die Sprechernotizen (oft enthalten die mehr Inhalt)

    Wir extrahieren beides, weil in den Notizen oft
    detailliertere Behauptungen stehen als auf der Folie.
    """
    prs = Presentation(file_path)
    sections = []
    full_text = ""

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = f"Folie {slide_num}:\n"

        # Text aus allen Textboxen
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"

        # Sprechernotizen
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                slide_text += f"\nNotizen: {notes}\n"

        sections.append({
            "heading": f"Folie {slide_num}",
            "text": slide_text
        })
        full_text += slide_text + "\n\n"

    return {
        "text": full_text,
        "sections": sections,
        "pages": len(prs.slides)
    }
